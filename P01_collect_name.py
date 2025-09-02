# collect_text_features_by_year.py
import os
import string
import pickle
import sys
import nltk
import pandas as pd

# 增加递归深度限制
sys.setrecursionlimit(2000)

# 下载 NLTK 的停用词数据
try:
    nltk.download('stopwords', quiet=True)
except Exception as e:
    print("下载 stopwords 出现错误：", e)

from nltk.corpus import stopwords
from docx import Document
from pptx import Presentation

# 预处理文本：小写、去标点、去停用词
def preprocess_text(text):
    text = text.lower()
    text = ''.join(char for char in text if char not in string.punctuation)
    stop_words = set(stopwords.words('english'))
    words = text.split()
    words = [word for word in words if word not in stop_words]
    return ' '.join(words)

# 提取各类型文件的前 len_doc 个字符
def extract_text_from_txt(file_path, len_doc):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read(len_doc)
        return content
    except Exception as e:
        print(f"读取 txt 文件 {file_path} 出错: {e}")
        return ""

def extract_text_from_docx(file_path, len_doc):
    try:
        doc = Document(file_path)
        full_text = [para.text for para in doc.paragraphs]
        content = "\n".join(full_text)[:len_doc]
        return content
    except Exception as e:
        print(f"读取 docx 文件 {file_path} 出错: {e}")
        return ""

def extract_text_from_pptx(file_path, len_doc):
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        content = "\n".join(full_text)[:len_doc]
        return content
    except Exception as e:
        print(f"读取 pptx 文件 {file_path} 出错: {e}")
        return ""

def extract_text_from_xlsx(file_path, len_doc):
    try:
        df_dict = pd.read_excel(file_path, sheet_name=None)
        texts = []
        for sheet_name, df in df_dict.items():
            texts.append(df.to_string())
        content = "\n".join(texts)[:len_doc]
        return content
    except Exception as e:
        print(f"读取 xlsx 文件 {file_path} 出错: {e}")
        return ""

# 核心特征提取函数
def extract_features_from_file(filepath, len_doc):
    ext = os.path.splitext(filepath)[1].lower()
    if ext == '.txt':
        content = extract_text_from_txt(filepath, len_doc)
    elif ext == '.docx':
        content = extract_text_from_docx(filepath, len_doc)
    elif ext == '.pptx':
        content = extract_text_from_pptx(filepath, len_doc)
    elif ext in ['.xls', '.xlsx']:
        content = extract_text_from_xlsx(filepath, len_doc)
    else:
        content = os.path.basename(filepath)
    content = content[:len_doc]
    return preprocess_text(content)

# 遍历目录获取所有文件路径
def get_all_file_paths(directories):
    file_paths = []
    for directory in directories:
        for entry in os.listdir(directory):
            path = os.path.join(directory, entry)
            if os.path.isfile(path):
                file_paths.append(path)
            elif os.path.isdir(path):
                file_paths.append(path)
                for sub_entry in os.listdir(path):
                    sub_path = os.path.join(path, sub_entry)
                    if os.path.isfile(sub_path) or os.path.isdir(sub_path):
                        file_paths.append(sub_path)
    return file_paths

# 按年份分别提取文本特征并保存到 file_embed_pkl_history 文件夹
def collect_data_and_features_by_year(directories_by_year, len_doc, save_folder):
    os.makedirs(save_folder, exist_ok=True)
    for year, directories in directories_by_year.items():
        file_paths = get_all_file_paths(directories)
        texts = []
        for file_path in file_paths:
            text = extract_features_from_file(file_path, len_doc)
            texts.append(text)
        save_file_str = os.path.join(save_folder, f'file_paths_and_texts_{year}.pkl')
        with open(save_file_str, 'wb') as f:
            pickle.dump((file_paths, texts), f)
        print(f"[{year}] 文件路径和文本特征已成功保存到 {save_file_str}")

# 标签生成与过滤
def generate_labels_and_filter(file_paths, texts, target_categories):
    filtered_file_paths = []
    filtered_texts = []
    labels = []
    normalized_target = set(os.path.normcase(os.path.normpath(cat)) for cat in target_categories)

    for path, text in zip(file_paths, texts):
        if path.lower().endswith('.lnk'):
            continue
        folder_name = os.path.basename(os.path.dirname(os.path.normpath(path)))
        folder_name = os.path.normcase(folder_name)

        filtered_file_paths.append(path)
        filtered_texts.append(text)
        if folder_name in normalized_target:
            labels.append(folder_name)
        else:
            labels.append("其他")
    return filtered_file_paths, filtered_texts, labels

# 拼接所有年份 pkl 并生成标签，保存最终数据集到 file_embed_pkl_history 文件夹
def combine_yearly_pkl_and_add_labels(years, target_categories, save_folder, final_save_file):
    all_file_paths = []
    all_texts = []

    for year in years:
        file_str = os.path.join(save_folder, f'file_paths_and_texts_{year}.pkl')
        with open(file_str, 'rb') as f:
            file_paths, texts = pickle.load(f)
            all_file_paths.extend(file_paths)
            all_texts.extend(texts)
        print(f"[{year}] 已加载 {len(file_paths)} 个文件")

    filtered_file_paths, filtered_texts, labels = generate_labels_and_filter(
        all_file_paths, all_texts, target_categories
    )

    final_save_path = os.path.join(save_folder, final_save_file)
    with open(final_save_path, 'wb') as f:
        pickle.dump((filtered_file_paths, filtered_texts, labels), f)
    print(f"最终拼接后的数据和标签已保存到 {final_save_path}")

# ========== 主流程 ========== #
if __name__ == "__main__":
    # 每年的目录

    """
        '2022': [r'E:\Alpha\StoreLatestYears\Store2022\M02广医事务性工作'],
        '2023': [r'E:\Alpha\StoreLatestYears\Store2023\M02广医事务性工作'],
        '2024': [r'E:\Alpha\StoreLatestYears\Store2024\M02广医事务性工作'],
    """


    directories_by_year = {
        '2025': [r'E:\Alpha\StoreLatestYears\Store2025\M02广医事务性工作']
    }

    # 目标分类列表
    target_categories = [
        "产学研_产业化工作",
        "产学研_社科科普",
        "产学研_科研工作",
        "人事工作_人才帽子",
        "人事工作_人才补贴政策",
        "人事工作_出境公开",
        "人事工作_提拔调动升职称",
        "人事工作_教师招聘或教师培训",
        "人事工作_职称评审",
        "外界公司",
        "学校党政行政_红头文件",
        "学院党政行政_红头文件",
        "工会后勤宣传保卫等工作",
        "年终绩效总结",
        "教学人才培养_教学学生竞赛",
        "教学人才培养_教改项目论文",
        "教学人才培养_本科教学",
        "教学人才培养_研究生教学",
        "日常整治与安全检查",
        "评审委员_校内校外",
        "财务_经费提醒"
    ]

    # 统一保存文件夹
    save_folder = 'file_embed_pkl_history'

    # 第一步：按年份分别提取并保存
    collect_data_and_features_by_year(directories_by_year, 200, save_folder)

    # 第二步：拼接所有年份并生成最终数据集
    final_save_file = 'file_paths_texts_and_labels_final.pkl'
    combine_yearly_pkl_and_add_labels(['2022', '2023', '2024', '2025'], target_categories, save_folder, final_save_file)
