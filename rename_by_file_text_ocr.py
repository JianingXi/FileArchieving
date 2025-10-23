# -*- coding: utf-8 -*-
"""
批量重命名工具：OCR + 文本提取 + PDF/PPT双模识别 + 关键词命名
依赖: pillow pytesseract jieba python-docx pymupdf python-pptx
"""
import os, re, csv, time, io
from pathlib import Path
from collections import Counter
import jieba
from collections import Counter
import pytesseract
from PIL import Image, ImageOps
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fitz  # PyMuPDF




IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}
TEXT_EXTS  = {".txt", ".docx", ".pdf", ".pptx", ".ppt"}
MAX_BASENAME_LEN = 80

# 设置Tesseract路径
for candidate in [
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe"
]:
    if os.path.exists(candidate):
        pytesseract.pytesseract.tesseract_cmd = candidate
        break

def sanitize_filename(name: str) -> str:
    # 先把英文停用词（如果混进来了）整体再删一次
    name = EN_STOPWORDS_RE.sub(" ", name)
    # 常规的非法字符与空白清理
    name = re.sub(r'[<>:"/\\|?*\x00-\x1F]', "_", name)
    name = re.sub(r"\s+", "_", name, flags=re.UNICODE).strip("_")
    return name or "ocr"














# —— 中文 & 英文停用词
STOPWORDS_CN = """
的 了 是 在 和 有 及 与 为 于 对 被 就 这 那 一个 一种 我 你 他 她 它 我们 他们 这些 那些 等 或 所以 因为 并 且 即 又 将 及其 以及 可 能 被 于 由
计算 可以 一切 开始 结束 输出 输入 
""".split()

# 更激进的英文停用词（含冠词、连词、介词、助动/情态动词、代词、副词等）
# 注意：我们同时构造一个基于 \b 的正则，在预清洗阶段直接把它们从文本里抹掉
STOPWORDS_EN_LIST = """
1 2 3 4 5 6 7 8 9 0
a an the and or but if while with within without from by on onto into over under
of in to at for about as per via during than through then so that which what who
whom whose it its is are was were be been being this that these those there here
where when why how not no nor do does did done doing can could may might must
shall should will would you your yours we our ours they them their theirs he she
him her his hers myself yourself itself ourselves yourselves themselves very also
just more most much many few some any all each either neither other another same
both ever never always often sometimes usually rather quite somewhat indeed
Fig data Data have has way people person
""".split()

# \b 边界 + IGNORECASE，确保 “The / THE / the” 都会被干净替换掉
EN_STOPWORDS_RE = re.compile(r"\b(?:%s)\b" % "|".join(map(re.escape, STOPWORDS_EN_LIST)), flags=re.IGNORECASE)

# 可保留的英文缩写（2~3 字母的有意义缩写，避免被长度过滤误杀）
WHITELIST_EN = {"AI","ML","ECG","CT","MRI","BME","NLP","KG","RL","GAN","EEG","EMG"}

def _preclean_text(text: str) -> str:
    """先对英文停用词做文本级‘抹除’，再统一压空白。"""
    text = EN_STOPWORDS_RE.sub(" ", text)
    text = re.sub(r"\s+", " ", text, flags=re.UNICODE).strip()
    return text

def extract_keywords(text: str, min_n: int = 5, max_n: int = 7):
    """
    中英文混合关键词提取（频次优先 + 首现位置；绝不保留英文虚词）。
    - 先把英文停用词从原文里删除（预清洗）
    - 中文：jieba 分词 + 中文停用词过滤 + 单字过滤
    - 英文：仅 A–Z 的词；长度 < 3 的一般丢弃（保留白名单缩写）
    - 合并后按 (频次 desc, 首现 index asc) 排序；去重后取 5–7 个
    """
    if not text:
        return ["ocr"]

    text = _preclean_text(text)

    # 中文词
    cn_words = [w for w in jieba.lcut(text)
                if len(w) > 1
                and w not in STOPWORDS_CN
                and re.search(r"[\u4e00-\u9fff]", w)]

    # 英文词（只允许字母）+ 长度过滤 + 白名单保留
    en_raw = re.findall(r"[A-Za-z]+", text)  # 只保留纯字母单词
    en_words = []
    for w in en_raw:
        wl = w.lower()
        # 再保险：即便预清洗了，也再做一次英文停用词拦截
        if EN_STOPWORDS_RE.fullmatch(w) or EN_STOPWORDS_RE.fullmatch(wl):
            continue
        if len(w) >= 3 or w.upper() in WHITELIST_EN:
            en_words.append(w)

    all_words = cn_words + en_words
    if not all_words:
        return ["ocr"]

    # 频次 + 首现位置（越靠前越优先）
    freq = Counter(all_words)
    sorted_words = sorted(freq.keys(), key=lambda w: (-freq[w], text.find(w)))

    # 去重并截断
    top, seen = [], set()
    for w in sorted_words:
        lw = w.lower()
        # 最终保险：绝不让英文虚词漏网
        if EN_STOPWORDS_RE.fullmatch(lw):
            continue
        if w not in seen:
            top.append(w)
            seen.add(w)
        if len(top) >= max_n:
            break

    # 不足 min_n 用剩余的非停用词补齐
    if len(top) < min_n:
        for w in sorted_words:
            lw = w.lower()
            if w not in seen and not EN_STOPWORDS_RE.fullmatch(lw):
                top.append(w); seen.add(w)
            if len(top) >= min_n:
                break

    return top[:max_n]









def open_image(path: Path):
    img = Image.open(path)
    if img.mode not in ("L", "RGB"):
        img = img.convert("RGB")
    gray = ImageOps.grayscale(img)
    gray = ImageOps.autocontrast(gray)
    return gray

def read_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except:
        return path.read_text(encoding="gbk", errors="ignore")

def read_docx_file(path: Path) -> str:
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def read_pdf_text_or_ocr(path):
    """PDF：先抽文本；若文本极少（< 40 字）则逐页渲染为图再 OCR。"""
    text = ""
    with fitz.open(path) as pdf:
        for page in pdf:
            text += page.get_text()
    if len(text.strip()) >= 40:
        return text

    # 无有效文本 → 逐页 OCR
    print(f"[INFO] PDF无文本，OCR识别: {path.name}")
    ocr_text = ""
    with fitz.open(path) as pdf:
        for page in pdf:
            pix = page.get_pixmap(dpi=220)  # 稍微高一点的 DPI 以提升 OCR 质量
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            gray = ImageOps.autocontrast(ImageOps.grayscale(img))
            ocr_text += pytesseract.image_to_string(gray, lang="chi_sim+eng") + "\n"
    return ocr_text

def read_pptx_text_or_ocr(path):
    """PPT/PPTX：先读取所有文本框；若几乎无文字，则提取每页里的图片逐个 OCR。"""
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text:
                texts.append(shape.text)
    merged = "\n".join(texts).strip()
    if len(merged) >= 20:
        return merged

    # 无有效文本 → OCR 幻灯片中的图片
    print(f"[INFO] PPT无文本，OCR识别图片: {path.name}")
    ocr_chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    pil_img = Image.open(io.BytesIO(shape.image.blob))
                    gray = ImageOps.autocontrast(ImageOps.grayscale(pil_img))
                    ocr_chunks.append(pytesseract.image_to_string(gray, lang="chi_sim+eng"))
                except Exception:
                    pass
    return "\n".join(ocr_chunks)

def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in IMAGE_EXTS:
        return pytesseract.image_to_string(open_image(path), lang="chi_sim+eng")
    elif ext == ".txt":
        return read_text_file(path)
    elif ext == ".docx":
        return read_docx_file(path)
    elif ext == ".pdf":
        return read_pdf_text_or_ocr(path)
    elif ext in (".pptx", ".ppt"):
        return read_pptx_text_or_ocr(path)
    return ""

def unique_path_in_dir(dirpath: Path, basename: str, ext: str) -> Path:
    candidate = dirpath / f"{basename}{ext}"
    i = 1
    while candidate.exists():
        candidate = dirpath / f"{basename}_{i}{ext}"
        i += 1
    return candidate

def should_process(p: Path):
    return p.suffix.lower() in IMAGE_EXTS.union(TEXT_EXTS)

def rename_files_with_keywords(root_dir: str):
    root_path = Path(root_dir)
    if not root_path.exists():
        print(f"[ERROR] 根目录不存在: {root_dir}")
        return
    log_path = root_path / f"rename_log_{time.strftime('%Y%m%d_%H%M%S')}.csv"
    with open(log_path, "w", newline="", encoding="utf-8-sig") as f:
        writer = csv.writer(f)
        writer.writerow(["original_path","new_path","extracted_keywords","text_preview"])
        for dirpath, _, filenames in os.walk(root_path):
            for fn in filenames:
                p = Path(dirpath) / fn
                if not should_process(p):
                    continue
                try:
                    text = extract_text(p)
                except Exception as e:
                    writer.writerow([p, "", "", f"READ_ERROR:{e}"])
                    continue
                keywords = extract_keywords(text, min_n=5, max_n=7)
                if not keywords:
                    keywords = ["ocr"]
                basename = sanitize_filename("_".join(keywords))[:MAX_BASENAME_LEN]
                dst = unique_path_in_dir(p.parent, basename, p.suffix.lower())
                try:
                    p.rename(dst)
                    print(f"[OK] {p.name} → {dst.name}")
                    writer.writerow([p, dst, "_".join(keywords), text[:120]])
                except Exception as e:
                    writer.writerow([p, "", "_".join(keywords), f"RENAME_ERROR:{e}"])
    print(f"完成，日志已保存：{log_path}")

# 调用示例
if __name__ == "__main__":
    ROOT_DIR = r"C:\Users\xijia\Desktop\新建文件夹"
    rename_files_with_keywords(ROOT_DIR)
