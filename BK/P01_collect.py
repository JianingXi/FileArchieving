import os
import string
import shutil
import pickle
import zipfile
import rarfile
import fitz  # PyMuPDF
import pandas as pd
from docx import Document
from pptx import Presentation
from nltk.corpus import stopwords
import nltk
import sys

# 增加递归深度限制
sys.setrecursionlimit(2000)

# 下载NLTK的停用词数据
nltk.download('stopwords')

# 文本预处理函数
def preprocess_text(text):
    # 转换为小写
    text = text.lower()
    # 去除所有标点符号
    text = ''.join(char for char in text if char not in string.punctuation)
    # 去除停用词
    stop_words = set(stopwords.words('english'))
    words = text.split()
    words = [word for word in words if word not in stop_words]
    return ' '.join(words)

# 从文件名中提取特征
def extract_features_from_filename(filename):
    filename = os.path.basename(filename).lower()  # 只获取文件名部分
    keywords = preprocess_text(filename)
    return keywords

# 从文本文件中提取内容
def extract_text_from_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        print(f"Failed to read txt file {file_path}: {e}")
        return ""

# 从Word文件中提取内容
def extract_text_from_doc(file_path):
    import win32com.client as win32
    try:
        word = win32.Dispatch('Word.Application')
        doc = word.Documents.Open(file_path)
        full_text = []
        for para in doc.Paragraphs:
            full_text.append(para.Range.Text)
        doc.Close()
        return '\n'.join(full_text)
    except Exception as e:
        print(f"Failed to read doc file {file_path}: {e}")
        return ""

def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        print(f"Failed to read docx file {file_path}: {e}")
        return ""

# 从PPT文件中提取内容
def extract_text_from_ppt(file_path):
    import win32com.client as win32
    try:
        powerpoint = win32.Dispatch('PowerPoint.Application')
        presentation = powerpoint.Presentations.Open(file_path)
        full_text = []
        for slide in presentation.Slides:
            for shape in slide.Shapes:
                if hasattr(shape, "TextFrame") and shape.TextFrame.HasText:
                    full_text.append(shape.TextFrame.TextRange.Text)
        presentation.Close()
        return '\n'.join(full_text)
    except Exception as e:
        print(f"Failed to read ppt file {file_path}: {e}")
        return ""

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return '\n'.join(full_text)
    except Exception as e:
        print(f"Failed to read pptx file {file_path}: {e}")
        return ""

# 从Excel文件中提取内容
def extract_text_from_xls(file_path):
    try:
        df = pd.read_excel(file_path, sheet_name=None)
        full_text = []
        for sheet_name, sheet_data in df.items():
            full_text.append(sheet_data.to_string())
        return '\n'.join(full_text)
    except Exception as e:
        print(f"Failed to read xls file {file_path}: {e}")
        return ""

def extract_text_from_xlsx(file_path):
    try:
        df = pd.read_excel(file_path, sheet_name=None)
        full_text = []
        for sheet_name, sheet_data in df.items():
            full_text.append(sheet_data.to_string())
        return '\n'.join(full_text)
    except Exception as e:
        print(f"Failed to read xlsx file {file_path}: {e}")
        return ""

# 从PDF文件中提取内容
def extract_text_from_pdf(file_path):
    try:
        doc = fitz.open(file_path)
        full_text = []
        for page in doc:
            full_text.append(page.get_text())
        return '\n'.join(full_text)
    except Exception as e:
        print(f"Failed to read pdf file {file_path}: {e}")
        return ""

# 从ZIP压缩包中提取内容
def extract_text_from_zip(file_path):
    temp_dir = 'temp_extracted_zip'
    os.makedirs(temp_dir, exist_ok=True)
    all_text = []
    try:
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        for root, dirs, files in os.walk(temp_dir):
            for file in files:
                file_path = os.path.join(root, file)
                all_text.append(extract_text_from_file(file_path))
    except Exception as e:
        print(f"Failed to read zip file {file_path}: {e}")
    finally:
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)
    return '\n'.join(all_text)

# 从RAR压缩包中提取内容
def extract_text_from_rar(file_path):
    temp_dir = 'temp_extracted_rar'
    os.makedirs(temp_dir, exist_ok=True)
    all_text = []
    try:
        with rarfile.RarFile(file_path) as rar_ref:
            rar_ref.extractall(temp_dir)
        for root, dirs, files in os.walk(temp_dir):
            for file in files:
                file_path = os.path.join(root, file)
                all_text.append(extract_text_from_file(file_path))
    except Exception as e:
        print(f"Failed to read rar file {file_path}: {e}")
    finally:
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)
    return '\n'.join(all_text)

# 根据文件类型提取内容
def extract_text_from_file(file_path):
    if file_path.endswith('.txt'):
        return extract_text_from_txt(file_path)
    elif file_path.endswith('.doc'):
        return extract_text_from_doc(file_path)
    elif file_path.endswith('.docx'):
        return extract_text_from_docx(file_path)
    elif file_path.endswith('.ppt'):
        return extract_text_from_ppt(file_path)
    elif file_path.endswith('.pptx'):
        return extract_text_from_pptx(file_path)
    elif file_path.endswith('.xls'):
        return extract_text_from_xls(file_path)
    elif file_path.endswith('.xlsx'):
        return extract_text_from_xlsx(file_path)
    elif file_path.endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith('.zip'):
        return extract_text_from_zip(file_path)
    elif file_path.endswith('.rar'):
        return extract_text_from_rar(file_path)
    else:
        return extract_features_from_filename(file_path)  # 默认提取文件名关键词

def get_all_file_paths(directories):
    file_paths = []
    for directory in directories:
        for entry in os.listdir(directory):
            path = os.path.join(directory, entry)
            if os.path.isfile(path):
                file_paths.append(path)
            elif os.path.isdir(path):
                file_paths.append(path)  # 将一级文件夹路径添加到列表中
                for sub_entry in os.listdir(path):  # 遍历一级文件夹中的文件和文件夹
                    sub_path = os.path.join(path, sub_entry)
                    if os.path.isfile(sub_path) or os.path.isdir(sub_path):
                        file_paths.append(sub_path)
    return file_paths

def collect_data_and_features(directories, save_file_str):
    # 获取所有文件和文件夹路径
    file_paths = get_all_file_paths(directories)

    texts = []
    for file_path in file_paths:
        text = extract_text_from_file(file_path)
        text = preprocess_text(text)
        texts.append(text)

    # 保存文件路径和文本到 .pkl 文件
    with open(save_file_str, 'wb') as f:
        pickle.dump((file_paths, texts), f)

    print("文件路径和文本特征已成功保存")

# 示例使用
directories = [r'D:\Alpha\M02广医事务性工作', r'D:\Alpha\StoreLatestYears\Store2023\M02广医事务性工作', r'D:\Alpha\StoreLatestYears\Store2024\M02广医事务性工作']  # 替换为你的文件夹路径，可以是单个地址或多个地址

collect_data_and_features(directories, '../file_paths_and_texts.pkl')
