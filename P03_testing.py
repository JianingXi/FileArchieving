import os
import string
import shutil
import pickle
import sys
import zipfile
import rarfile
import pandas as pd
import nltk
import numpy as np
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer

sys.setrecursionlimit(2000)

# 下载 NLTK 停用词
"""
try:
    nltk.download('stopwords', quiet=True)
except:
    pass
"""
from nltk.corpus import stopwords

def preprocess_text(text):
    text = text.lower()
    text = ''.join(c for c in text if c not in string.punctuation)
    stops = set(stopwords.words('english'))
    return ' '.join(w for w in text.split() if w not in stops)

def extract_text_from_txt(fp, L):
    try:
        with open(fp,'r',encoding='utf-8') as f:
            return f.read(L)
    except:
        return ""

def extract_text_from_docx(fp, L):
    try:
        doc = Document(fp)
        return "\n".join(p.text for p in doc.paragraphs)[:L]
    except:
        return ""

def extract_text_from_pptx(fp, L):
    try:
        prs = Presentation(fp)
        txt = []
        for s in prs.slides:
            for sh in s.shapes:
                if hasattr(sh, 'text'):
                    txt.append(sh.text)
        return "\n".join(txt)[:L]
    except:
        return ""

def extract_text_from_excel(fp, L):
    try:
        sheets = pd.read_excel(fp, sheet_name=None)
        return "\n".join(df.to_string() for df in sheets.values())[:L]
    except:
        return ""

def extract_text_from_zip(fp, L):
    try:
        with zipfile.ZipFile(fp) as z:
            names = z.namelist()
        return " ".join(preprocess_text(n) for n in names)[:L]
    except:
        return ""

def extract_text_from_rar(fp, L):
    try:
        with rarfile.RarFile(fp) as r:
            names = r.namelist()
        return " ".join(preprocess_text(n) for n in names)[:L]
    except:
        return ""

def extract_features_from_file(fp, L):
    ext = os.path.splitext(fp)[1].lower()
    if ext == '.txt':
        c = extract_text_from_txt(fp, L)
    elif ext == '.docx':
        c = extract_text_from_docx(fp, L)
    elif ext == '.pptx':
        c = extract_text_from_pptx(fp, L)
    elif ext in ('.xls', '.xlsx'):
        c = extract_text_from_excel(fp, L)
    elif ext == '.zip':
        c = extract_text_from_zip(fp, L)
    elif ext == '.rar':
        c = extract_text_from_rar(fp, L)
    else:
        c = ""
    name = os.path.basename(fp).lower().translate(str.maketrans('', '', string.punctuation))
    return (name + " " + preprocess_text(c))[:L]

def get_all_file_paths(dir_):
    paths = []
    for root, _, files in os.walk(dir_):
        for f in files:
            paths.append(os.path.join(root, f))
    return paths

def classify_files(src_dir, dst_root, L, thresh=0.5):
    # 获取所有文件
    fps = get_all_file_paths(src_dir)
    # 若目录为空，直接跳过
    if not fps:
        print(f"输入目录 '{src_dir}' 中没有文件，跳过分类。")
        return

    # 提取特征文本
    texts = [extract_features_from_file(fp, L) for fp in fps]

    # SBERT 嵌入
    embedder = SentenceTransformer('all-MiniLM-L6-v2')
    X = embedder.encode(texts, batch_size=32, show_progress_bar=False)

    # 加载模型和标签转换器
    with open('rf_model.pkl', 'rb') as f:
        model = pickle.load(f)
    with open('label_binarizer.pkl', 'rb') as f:
        mlb = pickle.load(f)
    labels = mlb.classes_.tolist()

    # 预测概率
    probas = model.predict_proba(X)
    P = np.vstack([p[:, 1] for p in probas]).T

    # 创建分类目录
    for lbl in labels:
        os.makedirs(os.path.join(dst_root, lbl), exist_ok=True)

    # 移动文件
    for fp, row in zip(fps, P):
        idxs = np.where(row >= thresh)[0].tolist()
        if not idxs:
            idxs = [int(np.argmax(row))]
        picked = [labels[i] for i in idxs]

        for lbl in picked:
            dst = os.path.join(dst_root, lbl, os.path.basename(fp))
            print(f"移动文件 {fp} 到 {dst} （分类：{lbl}）")
            shutil.move(fp, dst)

if __name__ == '__main__':
    src = r"C:\MyDocument\ToDoList\D20_ToDailyNotice"
    dst = os.path.join(os.getcwd(), "DoneFileArchived")
    os.makedirs(dst, exist_ok=True)
    classify_files(src, dst, L=50, thresh=0.5)
    print("分类并移动完毕。")



