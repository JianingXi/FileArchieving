import os
import string
import pickle
import sys
import nltk
import pandas as pd

# 增加递归深度限制
sys.setrecursionlimit(2000)

# 下载 NLTK 的停用词数据（第一次运行时会下载），若下载出错则捕获异常
try:
    nltk.download('stopwords', quiet=True)
except Exception as e:
    print("下载 stopwords 出现错误：", e)

from nltk.corpus import stopwords
from docx import Document
from pptx import Presentation

def preprocess_text(text):
    """
    对文本进行预处理：
      - 转换为小写
      - 去除所有标点符号
      - 去除英文停用词
    """
    text = text.lower()
    text = ''.join(char for char in text if char not in string.punctuation)
    stop_words = set(stopwords.words('english'))
    words = text.split()
    words = [word for word in words if word not in stop_words]
    return ' '.join(words)

def extract_text_from_txt(file_path, len_doc):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read(len_doc)  # 只读取前500个字符
        return content
    except Exception as e:
        print(f"读取 txt 文件 {file_path} 时出错: {e}")
        return ""

def extract_text_from_docx(file_path, len_doc):
    try:
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        content = "\n".join(full_text)[:len_doc]  # 只取前500个字符
        return content
    except Exception as e:
        print(f"读取 docx 文件 {file_path} 时出错: {e}")
        return ""

def extract_text_from_pptx(file_path, len_doc):
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        content = "\n".join(full_text)[:len_doc]
        return content
    except Exception as e:
        print(f"读取 pptx 文件 {file_path} 时出错: {e}")
        return ""

def extract_text_from_xlsx(file_path, len_doc):
    try:
        df_dict = pd.read_excel(file_path, sheet_name=None)
        texts = []
        for sheet_name, df in df_dict.items():
            texts.append(df.to_string())
        content = "\n".join(texts)[:len_doc]
        return content
    except Exception as e:
        print(f"读取 xlsx 文件 {file_path} 时出错: {e}")
        return ""

def extract_features_from_file(filepath, len_doc):
    """
    根据文件后缀提取文本特征：
      - 对于支持的文本型文件（txt, docx, pptx, xls/xlsx），读取内容的前500个字符；
      - 对于其他类型，则仅返回文件名。
    最后对提取的内容进行预处理。
    """
    ext = os.path.splitext(filepath)[1].lower()
    if ext == '.txt':
        content = extract_text_from_txt(filepath, len_doc)
    elif ext == '.docx':
        content = extract_text_from_docx(filepath, len_doc)
    elif ext == '.pptx':
        content = extract_text_from_pptx(filepath, len_doc)
    elif ext in ['.xls', '.xlsx']:
        content = extract_text_from_xlsx(filepath, len_doc)
    else:
        content = os.path.basename(filepath)
    # 仅保留前500个字符（如果内容较长）
    content = content[:len_doc]
    return preprocess_text(content)

def get_all_file_paths(directories):
    """
    遍历指定目录（支持多个目录），获取所有文件的完整路径，
    同时对子目录进行遍历。
    """
    file_paths = []
    for directory in directories:
        for entry in os.listdir(directory):
            path = os.path.join(directory, entry)
            if os.path.isfile(path):
                file_paths.append(path)
            elif os.path.isdir(path):
                file_paths.append(path)  # 如需保留文件夹路径也添加
                for sub_entry in os.listdir(path):
                    sub_path = os.path.join(path, sub_entry)
                    if os.path.isfile(sub_path) or os.path.isdir(sub_path):
                        file_paths.append(sub_path)
    return file_paths

def collect_data_and_features(directories, save_file_str, len_doc):
    """
    提取指定目录下所有文件的文本特征：
      - 对于支持的文本型文件（txt, docx, pptx, xls/xlsx），只读取前500个字符；
      - 对于其他文件，仅提取文件名；
    将所有文件路径和提取的文本特征保存到 pickle 文件中。
    """
    file_paths = get_all_file_paths(directories)
    texts = []
    for file_path in file_paths:
        text = extract_features_from_file(file_path, len_doc)
        texts.append(text)
    with open(save_file_str, 'wb') as f:
        pickle.dump((file_paths, texts), f)
    print("文件路径和文本特征已成功保存到", save_file_str)


# ----- 给文件打标签 ----- #
def generate_labels_and_filter(file_paths, texts, target_categories):
    """
    根据目标分类对每个文件路径生成标签：
      - 跳过所有 .lnk 文件
      - 如果文件所在文件夹的名称（归一化后）在 target_categories 中，
        则标签为该文件夹名称；
      - 否则标签设为 "其他"。
    同时过滤掉跳过的文件，并返回新的文件路径、文本和标签列表。
    """
    filtered_file_paths = []
    filtered_texts = []
    labels = []

    # 对目标分类进行归一化处理，用于比较
    normalized_target = set(os.path.normcase(os.path.normpath(cat)) for cat in target_categories)

    for path, text in zip(file_paths, texts):
        # 跳过 .lnk 文件
        if path.lower().endswith('.lnk'):
            continue
        # 提取最后一级文件夹名称，并归一化
        folder_name = os.path.basename(os.path.dirname(os.path.normpath(path)))
        folder_name = os.path.normcase(folder_name)

        filtered_file_paths.append(path)
        filtered_texts.append(text)
        # 如果文件所在文件夹在目标分类中，则标签为该文件夹名称，否则设为 "其他"
        if folder_name in normalized_target:
            labels.append(folder_name)
        else:
            labels.append("其他")
    return filtered_file_paths, filtered_texts, labels


def add_labels_to_features(features_file_str, target_categories, save_file_str):
    """
    从指定的 pickle 文件中加载数据（文件路径和文本，忽略已有标签），
    调用 generate_labels_and_filter 生成新的标签（并过滤 .lnk 文件），
    然后保存过滤后的数据到新的 pickle 文件中。
    """
    with open(features_file_str, 'rb') as f:
        data = pickle.load(f)

    if len(data) == 2:
        file_paths, texts = data
    elif len(data) == 3:
        file_paths, texts, _ = data
    else:
        raise ValueError("未知的数据格式，期望包含2或3个元素。")

    new_file_paths, new_texts, labels = generate_labels_and_filter(file_paths, texts, target_categories)

    with open(save_file_str, 'wb') as f:
        pickle.dump((new_file_paths, new_texts, labels), f)

    print("文件路径、文本和标签已成功保存到", save_file_str)











# 示例使用：请根据实际情况修改目录路径
directories = [
    r'D:\Alpha\StoreLatestYears\Store2022\M02广医事务性工作',
    r'D:\Alpha\StoreLatestYears\Store2023\M02广医事务性工作',
    r'D:\Alpha\StoreLatestYears\Store2024\M02广医事务性工作',
    r'D:\Alpha\StoreLatestYears\Store2025\M02广医事务性工作'
]

save_file_str = 'file_paths_and_texts.pkl'
collect_data_and_features(directories, save_file_str, 200)

# 标签：目标分类列表（这些字符串应与文件路径最后一级文件夹名称完全一致）
target_categories = [
    "产学研_产业化工作",
    "产学研_社科科普",
    "产学研_科研工作",
    "人事工作_人才帽子",
    "人事工作_人才补贴政策",
    "人事工作_出境公开",
    "人事工作_提拔调动升职称",
    "人事工作_教师招聘或教师培训",
    "人事工作_职称评审",
    "外界公司",
    "学校党政行政_红头文件",
    "学院党政行政_红头文件",
    "工会后勤宣传保卫等工作",
    "年终绩效总结",
    "教学人才培养_教学学生竞赛",
    "教学人才培养_教改项目论文",
    "教学人才培养_本科教学",
    "教学人才培养_研究生教学",
    "日常整治与安全检查",
    "评审委员_校内校外",
    "财务_经费提醒"
]

# 调用函数进行处理（请确保 'file_paths_texts_and_labels.pkl' 存在且格式正确）
add_labels_to_features('file_paths_and_texts.pkl', target_categories, 'file_paths_texts_and_labels.pkl')
